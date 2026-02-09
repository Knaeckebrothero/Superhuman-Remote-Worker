"""
Document Processor Core

Handles text extraction from various document formats (PDF, DOCX, PPTX, XLSX, TXT, HTML)
and applies intelligent chunking strategies for legal and technical documents.
"""

import re
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

from .document_models import (
    DocumentChunk,
    DocumentMetadata,
    DocumentCategory,
)

# Optional imports for document processing
try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    pdfplumber = None

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    DocxDocument = None

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False
    BeautifulSoup = None

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    load_workbook = None


# =============================================================================
# Chunking Configuration
# =============================================================================

CHUNKING_PRESETS = {
    "legal": {
        "max_chunk_size": 5000,
        "overlap_size": 500,
        "respect_boundaries": True,
        "boundary_patterns": [
            r"^(?:Article|Section|§|Artikel|Abschnitt)\s+\d+",  # Section headers
            r"^\(\d+\)",  # Numbered paragraphs
            r"^[a-z]\)",  # Lettered clauses
            r"^(?:\d+\.)+\s+",  # Numbered sections like "1.2.3 "
        ],
        "preserve_hierarchy": True,
    },
    "technical": {
        "max_chunk_size": 3500,
        "overlap_size": 350,
        "respect_boundaries": True,
        "boundary_patterns": [
            r"^#+\s+",  # Markdown headers
            r"^\d+\.\d+",  # Numbered sections
            r"^(?:REQ|FR|NFR)-\d+",  # Requirement IDs
        ],
        "preserve_hierarchy": True,
    },
    "general": {
        "max_chunk_size": 2000,
        "overlap_size": 200,
        "respect_boundaries": False,
        "boundary_patterns": [],
        "preserve_hierarchy": False,
    },
    "by_page": {
        "max_chunk_size": None,  # No limit - one chunk per page
        "overlap_size": 0,
        "respect_boundaries": True,
        "boundary_patterns": [r"\[PAGE \d+\]"],  # Page markers from PDF extraction
        "preserve_hierarchy": True,
        "page_based": True,
    },
}


# =============================================================================
# Language Detection
# =============================================================================

# Simple language detection based on common words
LANGUAGE_INDICATORS = {
    "de": [
        "und", "der", "die", "das", "ist", "sind", "werden", "muss", "soll",
        "gemäß", "sowie", "durch", "für", "bei", "zur", "vom", "nach", "über",
    ],
    "en": [
        "the", "and", "is", "are", "must", "shall", "should", "will", "may",
        "for", "with", "from", "this", "that", "which", "have", "has", "been",
    ],
}


def detect_language(text: str) -> str:
    """Simple language detection based on word frequency."""
    words = text.lower().split()[:500]  # Sample first 500 words
    word_set = set(words)

    scores = {}
    for lang, indicators in LANGUAGE_INDICATORS.items():
        score = sum(1 for w in indicators if w in word_set)
        scores[lang] = score

    return max(scores, key=scores.get) if scores else "unknown"


# =============================================================================
# Document Type Detection
# =============================================================================

DOCUMENT_TYPE_PATTERNS = {
    DocumentCategory.LEGAL: [
        r"\b(vertrag|contract|agreement|vereinbarung)\b",
        r"\b(gesetz|law|regulation|verordnung)\b",
        r"\b(§|artikel|article|section)\s+\d+",
        r"\b(GoBD|DSGVO|GDPR|HGB|AO)\b",
        r"\b(rechtlich|legal|binding|verbindlich)\b",
    ],
    DocumentCategory.TECHNICAL: [
        r"\b(API|SDK|interface|schnittstelle)\b",
        r"\b(system|component|module|komponente)\b",
        r"\b(specification|requirement|anforderung)\b",
        r"\b(architecture|design|implementation)\b",
    ],
    DocumentCategory.POLICY: [
        r"\b(policy|richtlinie|guideline|leitfaden)\b",
        r"\b(procedure|prozess|workflow|ablauf)\b",
        r"\b(standard|norm|compliance)\b",
    ],
}


def detect_document_type(text: str) -> DocumentCategory:
    """Detect document category based on content patterns."""
    text_lower = text.lower()

    scores = {cat: 0 for cat in DocumentCategory}

    for category, patterns in DOCUMENT_TYPE_PATTERNS.items():
        for pattern in patterns:
            matches = len(re.findall(pattern, text_lower, re.IGNORECASE))
            scores[category] += matches

    # Get category with highest score
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else DocumentCategory.GENERAL


# =============================================================================
# Text Extraction
# =============================================================================

class DocumentExtractor:
    """Extracts text content from various document formats."""

    def __init__(self):
        """Initialize the document extractor."""
        self._check_dependencies()

    # Text-based formats that can be read directly as plain text
    TEXT_FORMATS = {
        # Plain text
        ".txt", ".text",
        # Data formats
        ".csv", ".tsv",
        ".json", ".jsonl", ".ndjson",
        ".xml",
        ".yaml", ".yml",
        ".toml",
        ".ini", ".cfg", ".conf",
        # Markup/documentation
        ".md", ".markdown", ".rst", ".asciidoc", ".adoc",
        # Code files (common ones)
        ".py", ".js", ".ts", ".jsx", ".tsx",
        ".java", ".kt", ".scala",
        ".c", ".cpp", ".h", ".hpp",
        ".cs", ".go", ".rs", ".rb",
        ".php", ".swift", ".r",
        ".sql", ".sh", ".bash", ".zsh", ".ps1",
        ".css", ".scss", ".sass", ".less",
        # Config files
        ".env", ".properties", ".gradle",
        ".dockerfile",
        # Log files
        ".log",
        # Web
        ".svg",  # XML-based, readable as text
    }

    def _check_dependencies(self):
        """Log available document processing capabilities."""
        self.capabilities = {
            "pdf": PDF_AVAILABLE,
            "docx": DOCX_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "xlsx": XLSX_AVAILABLE,
            "html": HTML_AVAILABLE,
            "txt": True,  # Always available
            "text_formats": True,  # CSV, JSON, YAML, MD, code files, etc.
            "fallback": True,  # Attempt to read unknown formats as text
        }

    def extract(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text and metadata from a document.

        Args:
            file_path: Path to the document file

        Returns:
            Tuple of (extracted_text, extraction_info)

        Raises:
            ValueError: If file type is not supported or dependencies missing
            FileNotFoundError: If file does not exist
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = path.suffix.lower()

        if suffix == ".pdf":
            return self._extract_pdf(path)
        elif suffix == ".docx":
            return self._extract_docx(path)
        elif suffix == ".pptx":
            return self._extract_pptx(path)
        elif suffix == ".xlsx":
            return self._extract_xlsx(path)
        elif suffix == ".txt":
            return self._extract_txt(path)
        elif suffix in (".html", ".htm"):
            return self._extract_html(path)
        elif suffix in self.TEXT_FORMATS:
            # Known text-based formats - read as plain text
            return self._extract_txt(path)
        else:
            # Fallback: attempt to read as text
            return self._extract_fallback(path)

    def _extract_pdf(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF using pdfplumber."""
        if not PDF_AVAILABLE:
            raise ValueError(
                "PDF extraction requires pdfplumber. Install with: pip install pdfplumber"
            )

        text_parts = []
        page_count = 0
        info = {}

        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            info["page_count"] = page_count

            # Try to get metadata
            if pdf.metadata:
                info["title"] = pdf.metadata.get("Title", "")
                info["author"] = pdf.metadata.get("Author", "")

            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                if page_text:
                    # Add page marker for later reference
                    text_parts.append(f"[PAGE {page_num}]\n{page_text}")

        return "\n\n".join(text_parts), info

    def _extract_docx(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX using python-docx."""
        if not DOCX_AVAILABLE:
            raise ValueError(
                "DOCX extraction requires python-docx. Install with: pip install python-docx"
            )

        doc = DocxDocument(path)
        text_parts = []
        info = {}

        # Extract core properties
        if doc.core_properties:
            info["title"] = doc.core_properties.title or ""
            info["author"] = doc.core_properties.author or ""

        # Extract paragraphs with style information
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    text_parts.append(f"[H{level}] {para.text}")
                else:
                    text_parts.append(para.text)

        # Estimate page count (rough: ~3000 chars per page)
        total_text = "\n".join(text_parts)
        info["page_count"] = max(1, len(total_text) // 3000)

        return "\n\n".join(text_parts), info

    def _extract_txt(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from plain text file."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()

        # Estimate page count
        page_count = max(1, len(text) // 3000)

        return text, {"page_count": page_count}

    def _extract_fallback(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Attempt to read unknown file formats as text.

        This is a fallback for files with unrecognized extensions.
        It tries UTF-8 first, then falls back to latin-1 if that fails.
        Binary files (those with too many non-printable characters) will raise an error.
        """
        # First, check if file appears to be binary
        try:
            with open(path, "rb") as f:
                sample = f.read(8192)  # Read first 8KB to check

            # Check for null bytes (strong indicator of binary)
            if b"\x00" in sample:
                raise ValueError(
                    f"File appears to be binary: {path.suffix}. "
                    "Binary files cannot be processed as text."
                )

            # Check ratio of non-printable characters
            non_printable = sum(
                1 for b in sample
                if b < 32 and b not in (9, 10, 13)  # Allow tab, newline, carriage return
            )
            if len(sample) > 0 and non_printable / len(sample) > 0.1:
                raise ValueError(
                    f"File appears to be binary or encoded: {path.suffix}. "
                    f"Non-printable character ratio: {non_printable / len(sample):.1%}"
                )

        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"Could not read file {path}: {e}")

        # Try reading as text
        try:
            with open(path, "r", encoding="utf-8") as f:
                text = f.read()
        except UnicodeDecodeError:
            # Fallback to latin-1 which can read any byte sequence
            with open(path, "r", encoding="latin-1") as f:
                text = f.read()

        # Estimate page count
        page_count = max(1, len(text) // 3000)

        return text, {
            "page_count": page_count,
            "fallback_extraction": True,
            "original_extension": path.suffix,
        }

    def _extract_html(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from HTML using BeautifulSoup."""
        if not HTML_AVAILABLE:
            raise ValueError(
                "HTML extraction requires BeautifulSoup. Install with: pip install beautifulsoup4"
            )

        with open(path, "r", encoding="utf-8", errors="replace") as f:
            html = f.read()

        soup = BeautifulSoup(html, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer"]):
            element.decompose()

        # Get title
        title = ""
        if soup.title:
            title = soup.title.string or ""

        # Get text
        text = soup.get_text(separator="\n", strip=True)

        # Estimate page count
        page_count = max(1, len(text) // 3000)

        return text, {"page_count": page_count, "title": title}

    def _extract_pptx(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PowerPoint using python-pptx."""
        if not PPTX_AVAILABLE:
            raise ValueError(
                "PPTX extraction requires python-pptx. Install with: pip install python-pptx"
            )

        prs = Presentation(path)
        text_parts = []
        info = {}
        slide_count = len(prs.slides)
        info["page_count"] = slide_count

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                # Handle tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)

            if slide_texts:
                slide_content = "\n".join(slide_texts)
                text_parts.append(f"[SLIDE {slide_num}]\n{slide_content}")

        # Try to get title from first slide
        if prs.slides and prs.slides[0].shapes.title:
            info["title"] = prs.slides[0].shapes.title.text

        return "\n\n".join(text_parts), info

    def _extract_xlsx(self, path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Excel using openpyxl."""
        if not XLSX_AVAILABLE:
            raise ValueError(
                "XLSX extraction requires openpyxl. Install with: pip install openpyxl"
            )

        wb = load_workbook(path, data_only=True)
        text_parts = []
        info = {}
        info["page_count"] = len(wb.sheetnames)

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = [f"[SHEET: {sheet_name}]"]

            for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                # Skip completely empty rows
                if not any(cell is not None for cell in row):
                    continue

                # Format row as table-like structure
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if row_num == 1:
                    sheet_texts.append("| " + " | ".join(row_values) + " |")
                    sheet_texts.append("|" + "|".join(["---"] * len(row_values)) + "|")
                else:
                    sheet_texts.append("| " + " | ".join(row_values) + " |")

            if len(sheet_texts) > 1:  # More than just the header
                text_parts.append("\n".join(sheet_texts))

        wb.close()

        # Use first sheet name or filename as title
        info["title"] = wb.sheetnames[0] if wb.sheetnames else path.stem

        return "\n\n".join(text_parts), info


# =============================================================================
# Document Chunker
# =============================================================================

class DocumentChunker:
    """
    Intelligent document chunking with support for legal document structure.

    Handles:
    - Hierarchical document structure preservation
    - Boundary-respecting chunk splits
    - Overlap for context continuity
    - Token estimation for LLM compatibility
    """

    def __init__(
        self,
        max_chunk_size: Optional[int] = 1000,
        overlap_size: int = 200,
        respect_boundaries: bool = True,
        boundary_patterns: Optional[List[str]] = None,
        preserve_hierarchy: bool = True,
        page_based: bool = False,
    ):
        """
        Initialize the chunker.

        Args:
            max_chunk_size: Maximum characters per chunk (None for no limit)
            overlap_size: Characters to overlap between chunks
            respect_boundaries: Whether to split at natural boundaries
            boundary_patterns: Regex patterns for section boundaries
            preserve_hierarchy: Whether to track section hierarchy
            page_based: If True, create one chunk per page (ignores other settings)
        """
        self.max_chunk_size = max_chunk_size
        self.overlap_size = overlap_size
        self.respect_boundaries = respect_boundaries
        self.boundary_patterns = boundary_patterns or []
        self.preserve_hierarchy = preserve_hierarchy
        self.page_based = page_based

        # Compile boundary patterns
        self._compiled_patterns = [
            re.compile(p, re.MULTILINE) for p in self.boundary_patterns
        ]

    @classmethod
    def from_preset(cls, preset: str) -> "DocumentChunker":
        """Create a chunker from a preset configuration."""
        if preset not in CHUNKING_PRESETS:
            raise ValueError(f"Unknown preset: {preset}. Available: {list(CHUNKING_PRESETS.keys())}")

        config = CHUNKING_PRESETS[preset]
        return cls(**config)

    def chunk(self, text: str, document_id: str = "") -> List[DocumentChunk]:
        """
        Split document text into chunks.

        Args:
            text: Full document text
            document_id: Optional ID prefix for chunk IDs

        Returns:
            List of DocumentChunk objects
        """
        if not text.strip():
            return []

        # Detect structure first
        structure = self._detect_structure(text)

        # Perform chunking based on mode
        if self.page_based:
            chunks = self._chunk_by_page(text, structure)
        elif self.respect_boundaries:
            chunks = self._chunk_with_boundaries(text, structure)
        else:
            chunks = self._chunk_simple(text)

        # Build DocumentChunk objects
        doc_chunks = []
        for i, chunk_info in enumerate(chunks):
            chunk_id = f"{document_id}_chunk_{i:04d}" if document_id else f"chunk_{uuid.uuid4().hex[:8]}"

            # Calculate overlap with previous
            overlap = 0
            if i > 0:
                prev_end = chunks[i - 1]["end"]
                if chunk_info["start"] < prev_end:
                    overlap = prev_end - chunk_info["start"]

            doc_chunk = DocumentChunk(
                chunk_id=chunk_id,
                text=chunk_info["text"],
                start_position=chunk_info["start"],
                end_position=chunk_info["end"],
                section_hierarchy=chunk_info.get("hierarchy", []),
                chunk_index=i,
                overlap_with_previous=overlap,
                estimated_tokens=self._estimate_tokens(chunk_info["text"]),
                page_number=chunk_info.get("page"),
            )
            doc_chunks.append(doc_chunk)

        return doc_chunks

    def _detect_structure(self, text: str) -> Dict[str, Any]:
        """Detect document structure (sections, headings)."""
        structure = {
            "sections": [],
            "pages": [],
        }

        # Find page markers
        page_pattern = re.compile(r"\[PAGE (\d+)\]")
        for match in page_pattern.finditer(text):
            structure["pages"].append({
                "page": int(match.group(1)),
                "position": match.start(),
            })

        # Find section headers using boundary patterns
        for pattern in self._compiled_patterns:
            for match in pattern.finditer(text):
                structure["sections"].append({
                    "header": match.group().strip(),
                    "position": match.start(),
                })

        # Sort sections by position
        structure["sections"].sort(key=lambda x: x["position"])

        return structure

    def _chunk_by_page(self, text: str, structure: Dict[str, Any]) -> List[Dict]:
        """Create one chunk per page based on [PAGE N] markers.

        Args:
            text: Full document text with [PAGE N] markers
            structure: Structure info from _detect_structure

        Returns:
            List of chunk dictionaries, one per page
        """
        chunks = []
        pages = structure.get("pages", [])

        if not pages:
            # No page markers found - treat entire text as one chunk
            if text.strip():
                chunks.append({
                    "text": text.strip(),
                    "start": 0,
                    "end": len(text),
                    "hierarchy": [],
                    "page": 1,
                })
            return chunks

        # Add end marker for easier iteration
        pages_with_end = pages + [{"page": None, "position": len(text)}]

        for i in range(len(pages_with_end) - 1):
            page_info = pages_with_end[i]
            next_info = pages_with_end[i + 1]

            page_start = page_info["position"]
            page_end = next_info["position"]
            page_text = text[page_start:page_end]

            if page_text.strip():
                chunks.append({
                    "text": page_text.strip(),
                    "start": page_start,
                    "end": page_end,
                    "hierarchy": [],
                    "page": page_info["page"],
                })

        return chunks

    def _chunk_with_boundaries(self, text: str, structure: Dict[str, Any]) -> List[Dict]:
        """Chunk text respecting structural boundaries."""
        chunks = []
        sections = structure.get("sections", [])
        pages = structure.get("pages", [])

        # Add end marker
        sections_with_end = sections + [{"header": "", "position": len(text)}]

        current_hierarchy = []

        for i in range(len(sections_with_end) - 1):
            section_start = sections_with_end[i]["position"]
            section_end = sections_with_end[i + 1]["position"]
            section_text = text[section_start:section_end]
            section_header = sections_with_end[i]["header"]

            # Update hierarchy
            if self.preserve_hierarchy and section_header:
                current_hierarchy = self._update_hierarchy(current_hierarchy, section_header)

            # If section is small enough, add as single chunk
            if self.max_chunk_size is None or len(section_text) <= self.max_chunk_size:
                if section_text.strip():
                    chunks.append({
                        "text": section_text.strip(),
                        "start": section_start,
                        "end": section_end,
                        "hierarchy": current_hierarchy.copy(),
                        "page": self._find_page(section_start, pages),
                    })
            else:
                # Split section into smaller chunks
                sub_chunks = self._split_section(
                    section_text, section_start, current_hierarchy, pages
                )
                chunks.extend(sub_chunks)

        # Handle text before first section
        if sections and sections[0]["position"] > 0:
            pre_text = text[:sections[0]["position"]]
            if pre_text.strip():
                pre_chunks = self._split_section(pre_text, 0, [], pages)
                chunks = pre_chunks + chunks

        # If no sections found, use simple chunking
        if not chunks:
            chunks = self._chunk_simple(text)

        return chunks

    def _split_section(
        self, text: str, offset: int, hierarchy: List[str], pages: List[Dict]
    ) -> List[Dict]:
        """Split a section into smaller chunks with overlap."""
        chunks = []
        sentences = self._split_into_sentences(text)

        current_chunk = []
        current_len = 0
        chunk_start = offset

        for sentence in sentences:
            sentence_len = len(sentence)

            # If adding this sentence exceeds max, save current chunk
            if current_len + sentence_len > self.max_chunk_size and current_chunk:
                chunk_text = " ".join(current_chunk)
                chunks.append({
                    "text": chunk_text.strip(),
                    "start": chunk_start,
                    "end": chunk_start + len(chunk_text),
                    "hierarchy": hierarchy.copy(),
                    "page": self._find_page(chunk_start, pages),
                })

                # Start new chunk with overlap
                overlap_sentences = []
                overlap_len = 0
                for s in reversed(current_chunk):
                    if overlap_len + len(s) <= self.overlap_size:
                        overlap_sentences.insert(0, s)
                        overlap_len += len(s)
                    else:
                        break

                current_chunk = overlap_sentences
                current_len = overlap_len
                chunk_start = chunks[-1]["end"] - overlap_len

            current_chunk.append(sentence)
            current_len += sentence_len

        # Add final chunk
        if current_chunk:
            chunk_text = " ".join(current_chunk)
            chunks.append({
                "text": chunk_text.strip(),
                "start": chunk_start,
                "end": chunk_start + len(chunk_text),
                "hierarchy": hierarchy.copy(),
                "page": self._find_page(chunk_start, pages),
            })

        return chunks

    def _chunk_simple(self, text: str) -> List[Dict]:
        """Simple character-based chunking with overlap."""
        chunks = []
        start = 0

        while start < len(text):
            end = min(start + self.max_chunk_size, len(text))

            # Try to end at a sentence boundary
            if end < len(text):
                # Look for sentence end within last 100 chars
                search_start = max(end - 100, start)
                for punct in [". ", "! ", "? ", "\n\n"]:
                    last_punct = text.rfind(punct, search_start, end)
                    if last_punct > start:
                        end = last_punct + len(punct)
                        break

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "text": chunk_text,
                    "start": start,
                    "end": end,
                    "hierarchy": [],
                    "page": None,
                })

            # Move to next chunk with overlap
            start = end - self.overlap_size
            if start < 0:
                start = 0
            if start >= len(text):
                break

        return chunks

    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences."""
        # Simple sentence splitter
        sentence_endings = re.compile(r"(?<=[.!?])\s+")
        sentences = sentence_endings.split(text)
        return [s.strip() for s in sentences if s.strip()]

    def _update_hierarchy(self, current: List[str], new_header: str) -> List[str]:
        """Update section hierarchy based on new header."""
        # Simple hierarchy: just append for now
        # Could be enhanced to detect heading levels
        return current + [new_header]

    def _find_page(self, position: int, pages: List[Dict]) -> Optional[int]:
        """Find the page number for a given position."""
        if not pages:
            return None

        current_page = 1
        for page_info in pages:
            if page_info["position"] <= position:
                current_page = page_info["page"]
            else:
                break

        return current_page

    def _estimate_tokens(self, text: str) -> int:
        """Estimate token count (rough approximation: ~4 chars per token)."""
        return len(text) // 4


# =============================================================================
# Main Document Processor
# =============================================================================

class DocumentProcessor:
    """
    Main document processor combining extraction and chunking.

    Provides a high-level interface for processing documents
    through the ingestion pipeline.
    """

    def __init__(self, chunking_strategy: str = "legal"):
        """
        Initialize the document processor.

        Args:
            chunking_strategy: One of "legal", "technical", "general"
        """
        self.extractor = DocumentExtractor()
        self.chunker = DocumentChunker.from_preset(chunking_strategy)
        self.chunking_strategy = chunking_strategy

    def process(self, file_path: str) -> Tuple[DocumentMetadata, List[DocumentChunk]]:
        """
        Process a document file and return metadata and chunks.

        Args:
            file_path: Path to the document file

        Returns:
            Tuple of (DocumentMetadata, List[DocumentChunk])
        """
        path = Path(file_path)

        # Extract text
        text, info = self.extractor.extract(file_path)

        # Detect document properties
        language = detect_language(text)
        doc_category = detect_document_type(text)

        # Build metadata
        metadata = DocumentMetadata(
            title=info.get("title", path.stem),
            source=str(path.absolute()),
            document_type=doc_category,
            detected_language=language,
            page_count=info.get("page_count", 1),
            extraction_timestamp=datetime.utcnow(),
            jurisdiction=self._detect_jurisdiction(text, language),
            version=None,
            author=info.get("author"),
        )

        # Generate document ID for chunk naming
        doc_id = path.stem.replace(" ", "_")[:20]

        # Chunk the document
        chunks = self.chunker.chunk(text, document_id=doc_id)

        return metadata, chunks

    def _detect_jurisdiction(self, text: str, language: str) -> Optional[str]:
        """Attempt to detect jurisdiction from document content."""
        text_lower = text.lower()

        # German jurisdiction indicators
        if any(ind in text_lower for ind in ["bundesrepublik deutschland", "bgb", "hgb", "gobd"]):
            return "DE"

        # EU indicators
        if any(ind in text_lower for ind in ["european union", "eu regulation", "dsgvo", "gdpr"]):
            return "EU"

        # Default based on language
        if language == "de":
            return "DE"
        elif language == "en":
            return None  # Could be US, UK, etc.

        return None

    def get_capabilities(self) -> Dict[str, bool]:
        """Return available document processing capabilities."""
        return self.extractor.capabilities


# =============================================================================
# Utility Functions
# =============================================================================

def estimate_processing_time(file_path: str) -> float:
    """Estimate processing time in seconds based on file size."""
    path = Path(file_path)
    if not path.exists():
        return 0.0

    size_bytes = path.stat().st_size
    size_mb = size_bytes / (1024 * 1024)

    # Rough estimates: ~1 second per MB for PDF, ~0.5s for text
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return size_mb * 1.0
    elif suffix == ".docx":
        return size_mb * 0.8
    elif suffix == ".pptx":
        return size_mb * 1.2
    elif suffix == ".xlsx":
        return size_mb * 0.6
    else:
        return size_mb * 0.5
