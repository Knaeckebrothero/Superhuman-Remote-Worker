"""File operations tools for the Universal Agent.

Provides core file read/write/edit operations within the workspace.
These are the most commonly used tools for interacting with workspace files.

Enhanced with visual content support:
- Multimodal models receive rendered page screenshots (base64)
- Text-only models receive AI-generated descriptions of visual content
- Configurable via `llm.multimodal` in agent config
"""

import base64
import logging
import mimetypes
import zipfile
from pathlib import Path, PurePosixPath
from typing import Any, List, Literal, Optional

from langchain_core.tools import tool
from src.core.workspace_backend import WorkspaceUnavailableError

from src.services.image_content import (
    IMAGE_DATA_TAG_TEMPLATE,
    PAGE_IMAGE_TAG_TEMPLATE,
)
from src.services.cloud_mount.guardrails import workspace_path_touches_cloud
from src.utils.pdf import PDFReader, format_read_info
from ..context import ToolContext

from src.shared.tool_catalog.definitions import (
    FILE_TOOLS_METADATA as FILE_TOOLS_METADATA,
)

logger = logging.getLogger(__name__)

# Supported image file extensions
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}

# Supported audio file extensions (matching Whisper API supported formats)
AUDIO_EXTENSIONS = {
    ".mp3",
    ".wav",
    ".m4a",
    ".ogg",
    ".flac",
    ".webm",
    ".mp4",
    ".mpeg",
    ".mpga",
    ".oga",
    ".opus",
}

# Document extensions that support visual rendering
VISUAL_DOCUMENT_EXTENSIONS = {".pdf", ".pptx", ".docx"}

# Archive extensions read_file lists entries for instead of attempting to
# decode as text. Scoped to zip for now — the only format the session/worker
# upload paths extract (knowledge-base/knowledge/issues/session_uploads_never_extract_archives.md).
ARCHIVE_EXTENSIONS = {".zip"}

# Bound the entry listing for a large archive — read_file's job is to tell
# the agent what's inside, not to reproduce a multi-thousand-entry manifest.
MAX_ARCHIVE_LISTING_ENTRIES = 200

# Sidecar suffix the session-upload seam
# (orchestrator/services/thread_uploads.py::ZIP_REFUSAL_NOTE_SUFFIX) writes
# next to a zip it fell back to storing verbatim because extraction was
# refused (corrupt, traversal entry, over a cap). Checked before describing
# an archive below — otherwise a merely *parseable* but refused zip reads as
# an ordinary, successfully extracted one: a full entry listing with nothing
# saying those entries were never separated into readable files. That is
# the motivating incident's dead end, recreated one layer down. No shared
# import exists between the orchestrator and agent processes; keep both
# ends of this string in sync by hand.
ZIP_REFUSAL_NOTE_SUFFIX = ".extraction-refused.txt"

# Tool metadata for registry
# Phase availability: file tools are available in both strategic and tactical modes


def _cloud_cache_guard_for_path(context: ToolContext, path: str) -> Optional[str]:
    cloud_mount_cfg = context.get_config("cloud_mount", {})
    if not isinstance(cloud_mount_cfg, dict) or not cloud_mount_cfg.get("active"):
        return None
    if not workspace_path_touches_cloud(path):
        return None
    manager = cloud_mount_cfg.get("_manager")
    if manager is None or not hasattr(manager, "cache_limit_message"):
        return None
    try:
        return manager.cache_limit_message()
    except Exception as exc:
        logger.warning("Cloud cache guard check failed: %s", exc)
        return None


def _cloud_upperdir_guard_for_path(context: ToolContext, path: str) -> Optional[str]:
    """WRITE-scoped sibling of ``_cloud_cache_guard_for_path``: blocks writes

    that would copy-up into an over-quota overlay upperdir (design §7/§9.9).
    Reads never copy-up, so this is only wired into write_file/edit_file.
    """
    cloud_mount_cfg = context.get_config("cloud_mount", {})
    if not isinstance(cloud_mount_cfg, dict) or not cloud_mount_cfg.get("active"):
        return None
    if not workspace_path_touches_cloud(path):
        return None
    overlay = cloud_mount_cfg.get("_overlay_manager")
    if overlay is None or not hasattr(overlay, "quota_guard_message"):
        return None
    try:
        return overlay.quota_guard_message()
    except Exception as exc:
        logger.warning("Cloud upperdir guard check failed: %s", exc)
        return None


def _absolute(workspace: Any, path: str) -> str:
    """Render ``path`` as the absolute path the shell would see.

    The file tools resolve relative paths against the **workspace root**;
    ``shell_execute`` resolves them against the tab's current directory, which
    any ``cd`` moves. So the same relative string can denote two different
    files, and echoing the caller's own string back in a tool result tells the
    model nothing about which one it got. Job bbce4bed lost a deliverable in
    exactly that gap: ``Written: output/report.md`` next to a shell that was
    ``cd``'d into a subdirectory, with nothing anywhere naming the anchor.

    Delegates to the backend's own ``resolve_path`` (via
    :meth:`WorkspaceManager.get_path`), the same resolution ``_file_not_found``
    reports, so every backend — remote, subdir, virtual overlay — renders its
    real root. Falls back to the relative path if resolution fails: a
    diagnostic aid must never be the reason a write fails.

    See knowledge-base/knowledge/issues/deliverable_lost_to_nested_repo_commit_and_stranded_mode_a_job.md.
    """
    try:
        return str(workspace.get_path(path))
    except Exception:  # pragma: no cover - defensive; write already succeeded
        return path


def _get_mime_type(file_path: Path) -> str:
    """Get MIME type for a file based on extension."""
    mime_type, _ = mimetypes.guess_type(str(file_path))
    return mime_type or "application/octet-stream"


def _is_image_file(file_path: Path) -> bool:
    """Check if file is a supported image format."""
    return file_path.suffix.lower() in IMAGE_EXTENSIONS


def _is_audio_file(file_path: Path) -> bool:
    """Check if file is a supported audio format."""
    return file_path.suffix.lower() in AUDIO_EXTENSIONS


def _is_visual_document(file_path: Path) -> bool:
    """Check if file is a document that supports visual rendering."""
    return file_path.suffix.lower() in VISUAL_DOCUMENT_EXTENSIONS


def _is_archive_file(file_path: Path) -> bool:
    """Check if file is an archive read_file should list rather than decode."""
    return file_path.suffix.lower() in ARCHIVE_EXTENSIONS


def _read_zip_extraction_note(workspace: Any, path: str) -> Optional[str]:
    """The session-upload seam's refusal note for ``path``, if one exists.

    Returns None for the overwhelming majority of archives — either the
    upload extracted cleanly, this zip predates the extraction seam, or it
    was never a session upload at all (e.g. one living inside a cloned
    repo). Never raises: a missing or unreadable note must never block the
    ordinary archive listing that follows it.
    """
    note_path = f"{path}{ZIP_REFUSAL_NOTE_SUFFIX}"
    try:
        if not workspace.exists(note_path):
            return None
        return workspace.read_file(note_path).strip()
    except Exception as e:
        logger.debug(f"Could not read zip extraction note for {path}: {e}")
        return None


def _describe_zip_archive(local_path: Path, display_name: str, size: int) -> str:
    """Build an entry listing for a zip archive instead of its raw bytes.

    Falls back to the generic binary message when ``local_path`` wears a
    ``.zip`` extension but isn't actually a readable archive — a corrupt
    upload or a renamed non-zip file fails cleanly either way, never with a
    stack trace or a raw codec error.
    """
    try:
        with zipfile.ZipFile(local_path) as zf:
            infos = [
                info
                for info in zf.infolist()
                if not info.is_dir()
                and "__MACOSX" not in PurePosixPath(info.filename).parts
                and not any(
                    part.startswith(".") for part in PurePosixPath(info.filename).parts
                )
            ]
    except Exception as e:
        logger.debug(f"Could not list zip entries for {display_name}: {e}")
        return f"[binary file: {display_name}, {size:,} bytes]"

    lines = [f"Archive: {display_name} — {len(infos)} file(s), {size:,} bytes"]
    lines.append("")
    shown = infos[:MAX_ARCHIVE_LISTING_ENTRIES]
    for info in shown:
        lines.append(f"{info.file_size:>12,}  {info.filename}")
    remaining = len(infos) - len(shown)
    if remaining > 0:
        lines.append(f"… and {remaining} more")
    lines.append("")
    lines.append(
        "This is a zip archive — read_file lists its entries but cannot "
        "show file contents directly."
    )
    return "\n".join(lines)


def create_file_tools(context: ToolContext) -> List[Any]:
    """Create file operation tools with injected context.

    Args:
        context: ToolContext with workspace_manager

    Returns:
        List of LangChain tool functions

    Raises:
        ValueError: If context doesn't have a workspace_manager
    """
    if not context.has_workspace():
        raise ValueError("ToolContext must have a workspace_manager for file tools")

    workspace = context.workspace_manager

    def _file_not_found(path: str) -> str:
        resolved_path = workspace.get_path(path)
        return (
            f"Error: File not found: {resolved_path}\n"
            f'  (resolved from workspace root; you passed "{path}")\n'
            "  Use `search_files` to locate it if you expected it elsewhere."
        )

    # Get word limit (with backward compatibility fallback)
    max_read_words = context.get_config("max_read_words")
    if max_read_words is None:
        # Fall back to legacy bytes limit, convert to words
        max_read_size_legacy = context.get_config(
            "max_read_size", 137_500
        )  # ~25k words
        max_read_words = int(max_read_size_legacy / 5.5)

    # Token-derived ceiling (session_silent_failure_audit.md #5): one tool
    # result must not occupy more than ~15% of the main model's context
    # window (~0.75 words/token). The configured max_read_words still rules
    # when smaller; this bites only when the window is small relative to it —
    # four unbounded PDF reads once filled a 128k window to 183%.
    model_window = context.get_config("model_max_context_tokens")
    if model_window:
        derived_cap = int(model_window * 0.15 * 0.75)
        if derived_cap < max_read_words:
            max_read_words = max(derived_cap, 1_000)

    # Initialize PDF reader with word limit
    pdf_reader = PDFReader(max_words_per_read=max_read_words)

    # Line-based reading constants (matching Claude Code behavior)
    DEFAULT_LINE_LIMIT = 2000
    MAX_LINE_LIMIT = 2000
    MAX_LINE_LENGTH = 2000  # Truncate lines longer than this

    def _read_pdf_file(
        full_path,
        relative_path: str,
        page_start: Optional[int],
        page_end: Optional[int],
    ) -> str:
        """Internal helper to read PDF files with page support."""
        if not pdf_reader.is_available():
            return "Error: PDF reading requires pdfplumber. Install with: pip install pdfplumber"

        try:
            text, read_info = pdf_reader.read_pages(
                full_path, page_start=page_start, page_end=page_end
            )

            # Build header showing what was read
            pages = read_info["pages_read"]
            if len(pages) == 1:
                header = f"[Page {pages[0]} of {read_info['total_pages']}]"
            elif len(pages) > 1:
                header = f"[Pages {pages[0]}-{pages[-1]} of {read_info['total_pages']}]"
            else:
                header = f"[No pages read - total pages: {read_info['total_pages']}]"

            # Build result
            result_parts = [header, "", text]

            # Add continuation guidance if truncated
            if read_info["was_truncated"]:
                result_parts.append("")
                result_parts.append(format_read_info(read_info, relative_path))

            return "\n".join(result_parts)

        except ValueError as e:
            return f"Error: {str(e)}"
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"PDF read error for {relative_path}: {e}")
            return f"Error reading PDF: {str(e)}"

    def _handle_image_file(
        local_path: Path,
        describe: Optional[str],
        display_name: str = "",
    ) -> str:
        """Handle standalone image files.

        For multimodal models: Returns base64-encoded image data.
        For text-only models: Returns AI-generated description.

        Args:
            local_path: Path on the local filesystem (may be a temp file
                        downloaded from a remote workspace).
            describe: Optional query for visual analysis.
            display_name: Original filename for display headers.
        """
        name = display_name or local_path.name
        if context.get_phase_multimodal():
            # Return image for multimodal model to see directly
            try:
                image_data = local_path.read_bytes()
                base64_image = base64.b64encode(image_data).decode()
                mime_type = _get_mime_type(local_path)

                # The tag is stripped + replaced with a marker by the
                # graph-side post-processor (`extract_image_tags` in
                # `src/services/image_content.py`), which also appends a
                # synthesized HumanMessage carrying the image as a real
                # provider content block. Multimodal primary models see
                # the actual image; the cleaned ToolMessage stays small.
                return (
                    f"[IMAGE: {name}]\n"
                    f"Type: {mime_type}\n"
                    f"Size: {len(image_data):,} bytes\n\n"
                    + IMAGE_DATA_TAG_TEMPLATE.format(mime=mime_type, b64=base64_image)
                )
            except WorkspaceUnavailableError:
                raise
            except Exception as e:
                logger.error(f"Error reading image {local_path}: {e}")
                return f"Error reading image: {str(e)}"
        else:
            # Get AI description for text-only model
            try:
                from src.services.vision_helper import get_vision_helper
                from src.services.description_cache import get_description_cache

                cache = get_description_cache()
                vision = get_vision_helper()

                # Check cache first
                cached = cache.get(local_path, query=describe)
                if cached:
                    logger.debug(f"Cache hit for image: {name}")
                    return f"[IMAGE: {name}]\n\n{cached}"

                # Generate description
                image_data = local_path.read_bytes()
                description = vision.describe_image_sync(
                    image_data,
                    mime_type=_get_mime_type(local_path),
                    query=describe,
                    job_id=context.job_id,
                )

                # Cache for future use
                cache.set(local_path, description, query=describe)

                return f"[IMAGE: {name}]\n\n{description}"

            except ImportError as e:
                logger.warning(f"Vision services not available: {e}")
                return (
                    f"[IMAGE: {name}]\n"
                    f"(Visual description not available - vision services not configured)"
                )
            except WorkspaceUnavailableError:
                raise
            except Exception as e:
                logger.error(f"Error describing image {local_path}: {e}")
                return f"[IMAGE: {name}]\n(Error generating description: {str(e)})"

    def _handle_audio_file(
        local_path: Path,
        offset: Optional[int] = None,
        limit: Optional[int] = None,
        display_name: str = "",
    ) -> str:
        """Handle audio files via Whisper transcription with line-numbered paging.

        Transcribes audio to text using AudioHelper (large files are
        transparently chunked). Results are cached via DescriptionCache.
        Output is line-numbered with offset/limit paging, matching text files.

        Args:
            local_path: Path on the local filesystem (may be a temp file
                        downloaded from a remote workspace).
            offset: Starting line number (1-indexed).
            limit: Number of lines to return.
            display_name: Original filename for display headers.
        """
        name = display_name or local_path.name
        try:
            from src.services.audio_helper import (
                get_audio_helper,
                split_transcript_into_lines,
            )
            from src.services.description_cache import get_description_cache

            cache = get_description_cache()

            # Get transcript from cache or via transcription
            cached = cache.get(local_path)
            if cached:
                logger.debug(f"Cache hit for audio: {name}")
                transcript = cached
            else:
                audio = get_audio_helper()
                transcript = audio.transcribe_sync(
                    local_path,
                    job_id=context.job_id,
                )

                # Don't cache error results
                if not transcript.startswith("[Error"):
                    cache.set(local_path, transcript)

            # Return errors directly (no line numbering)
            if transcript.startswith("[Error"):
                return f"[AUDIO: {name}]\n\n{transcript}"

            # Split into lines for paging
            lines = split_transcript_into_lines(transcript)
            total_lines = len(lines)

            if total_lines == 0:
                return f"[AUDIO: {name}]\n\n(No speech content detected)"

            # Apply offset/limit defaults (same as text files)
            start_line = offset if offset is not None else 1
            line_count = limit if limit is not None else DEFAULT_LINE_LIMIT
            line_count = min(line_count, MAX_LINE_LIMIT)

            if start_line < 1:
                return "Error: offset must be >= 1 (line numbers are 1-indexed)"

            if start_line > total_lines:
                return (
                    f"Error: offset ({start_line}) exceeds total lines ({total_lines})"
                )

            # Extract requested range
            end_line = min(start_line + line_count - 1, total_lines)
            selected_lines = lines[start_line - 1 : end_line]

            # Format with line numbers (cat -n style)
            output_lines = []
            for i, line in enumerate(selected_lines, start=start_line):
                if len(line) > MAX_LINE_LENGTH:
                    line = line[:MAX_LINE_LENGTH] + "..."
                output_lines.append(f"{i:6}\t{line}")

            result = f"[AUDIO: {name}]\n\n" + "\n".join(output_lines)

            # Word count cap (same as text files)
            word_count = len(result.split())
            if word_count > max_read_words:
                words = result.split()
                truncated_text = " ".join(words[:max_read_words])
                last_newline = truncated_text.rfind("\n")
                if last_newline > 0:
                    truncated_text = truncated_text[:last_newline]
                kept_lines = truncated_text.count("\n")  # -1 for header line
                actual_end = start_line + kept_lines - 1
                result = truncated_text
                result += (
                    f"\n\n[TRUNCATED at word limit ({max_read_words:,} words). "
                    f"Showing lines {start_line}-{actual_end} of {total_lines}. "
                    f"Use offset={actual_end + 1} to continue.]"
                )
            elif end_line < total_lines:
                result += (
                    f"\n\n[Lines {start_line}-{end_line} of {total_lines}. "
                    f"Use offset={end_line + 1} to continue.]"
                )

            return result

        except ImportError as e:
            logger.warning(f"Audio services not available: {e}")
            return (
                f"[AUDIO: {name}]\n"
                f"(Transcription not available - audio services not configured)"
            )
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"Error transcribing audio {local_path}: {e}")
            return f"[AUDIO: {name}]\n(Error transcribing: {str(e)})"

    def _get_visual_content(
        full_path: Path,
        page_num: int,
        describe: Optional[str],
    ) -> str:
        """Get visual content description for a document page.

        For multimodal models: Returns base64-encoded page screenshot.
        For text-only models: Returns AI-generated description.
        """
        try:
            from src.services.document_renderer import get_document_renderer
            from src.services.vision_helper import get_vision_helper
            from src.services.description_cache import get_description_cache

            renderer = get_document_renderer()

            # Render the page as PNG. Per-family DPI (matrix settings.pdf_render_dpi
            # via limits) lets patch-model mains render fewer pixels than the
            # provider downscales away; None -> renderer default (150).
            try:
                page_image = renderer.render_page(
                    full_path, page_num, dpi=context.get_config("pdf_render_dpi")
                )
            except WorkspaceUnavailableError:
                raise
            except Exception as e:
                logger.warning(
                    f"Could not render page {page_num} of {full_path.name}: {e}"
                )
                return ""  # No visual content available

            if context.get_phase_multimodal():
                # Tag is stripped + delivered as a real image content block
                # by the graph-side post-processor; see
                # `src/services/image_content.py`.
                base64_image = base64.b64encode(page_image).decode()
                return "\n" + PAGE_IMAGE_TAG_TEMPLATE.format(
                    page=page_num, mime="image/png", b64=base64_image
                )
            else:
                # Get AI description for text-only model
                cache = get_description_cache()
                vision = get_vision_helper()

                # Check cache first
                cached = cache.get(full_path, page=page_num, query=describe)
                if cached:
                    logger.debug(f"Cache hit for page {page_num} of {full_path.name}")
                    return f"\n[PAGE {page_num} - VISUAL CONTENT]\n{cached}"

                # Generate description
                description = vision.describe_document_page_sync(
                    page_image,
                    page_num=page_num,
                    query=describe,
                    job_id=context.job_id,
                )

                # Cache for future use
                cache.set(full_path, description, page=page_num, query=describe)

                if describe:
                    return f'\n[PAGE {page_num} - VISUAL CONTENT (Query: "{describe[:50]}...")]\n{description}'
                else:
                    return f"\n[PAGE {page_num} - VISUAL CONTENT]\n{description}"

        except ImportError as e:
            logger.debug(f"Vision services not available: {e}")
            return ""  # Silently skip visual content if services not available
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.warning(f"Error getting visual content for page {page_num}: {e}")
            return ""

    def _read_visual_document(
        full_path: Path,
        relative_path: str,
        page_start: Optional[int],
        page_end: Optional[int],
        describe: Optional[str],
    ) -> str:
        """Read a document with visual content (PDF, PPTX, DOCX).

        Combines text extraction with visual content based on multimodal setting.
        """
        suffix = full_path.suffix.lower()

        # For PDF, use existing reader for text
        if suffix == ".pdf":
            text_result = _read_pdf_file(full_path, relative_path, page_start, page_end)

            # If there was an error, return it
            if text_result.startswith("Error:"):
                return text_result

            # Get page range that was read
            try:
                from src.services.document_renderer import get_document_renderer

                renderer = get_document_renderer()
                total_pages = renderer.get_page_count(full_path)

                start = page_start or 1
                end = min(page_end or total_pages, total_pages)

                # Decide which pages actually need rasterizing. Text-rich,
                # image-free pages are already represented by the extracted
                # text above, so rendering them only burns image tokens
                # (context_token_accounting.md S2). Fail-open: an empty map
                # (inspection failed) renders every page, exactly as before.
                from src.utils.pdf import compress_ranges, page_render_decisions

                decisions = page_render_decisions(full_path, start, end)

                visual_parts = []
                skipped = []
                for page_num in range(start, end + 1):
                    decision = decisions.get(page_num)
                    if decision is not None and not decision["render"]:
                        skipped.append(page_num)
                        continue
                    visual_content = _get_visual_content(full_path, page_num, describe)
                    if visual_content:
                        visual_parts.append(visual_content)

                parts = [text_result]
                if visual_parts:
                    parts.append("\n".join(visual_parts))
                if skipped:
                    parts.append(
                        f"\n[Did not rasterize {len(skipped)} text-only "
                        f"page(s) — already included as text above: "
                        f"pages {compress_ranges(skipped)}]"
                    )
                return "\n".join(parts)

            except WorkspaceUnavailableError:
                raise
            except Exception as e:
                logger.debug(f"Could not add visual content: {e}")
                return text_result

        # For PPTX and DOCX, we need different text extraction
        elif suffix == ".pptx":
            return _read_pptx_file(
                full_path, relative_path, page_start, page_end, describe
            )
        elif suffix == ".docx":
            return _read_docx_file(
                full_path, relative_path, page_start, page_end, describe
            )
        else:
            return f"Error: Unsupported visual document type: {suffix}"

    def _read_pptx_file(
        full_path: Path,
        relative_path: str,
        slide_start: Optional[int],
        slide_end: Optional[int],
        describe: Optional[str],
    ) -> str:
        """Read a PowerPoint file with text and visual content."""
        try:
            from pptx import Presentation

            prs = Presentation(full_path)
            total_slides = len(prs.slides)

            start = slide_start or 1
            end = min(slide_end or total_slides, total_slides)

            if start > total_slides:
                return f"Error: slide_start ({start}) exceeds total slides ({total_slides})"

            result_parts = [f"[Slides {start}-{end} of {total_slides}]", ""]

            for slide_num in range(start, end + 1):
                slide = prs.slides[slide_num - 1]

                # Extract text from slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                result_parts.append(f"[SLIDE {slide_num}]")
                if slide_text:
                    result_parts.append("\n".join(slide_text))
                else:
                    result_parts.append("(No text content)")

                # Add visual content
                visual_content = _get_visual_content(full_path, slide_num, describe)
                if visual_content:
                    result_parts.append(visual_content)

                result_parts.append("")  # Blank line between slides

            return "\n".join(result_parts)

        except ImportError:
            return "Error: python-pptx not installed. Install with: pip install python-pptx"
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"PPTX read error for {relative_path}: {e}")
            return f"Error reading PowerPoint: {str(e)}"

    def _read_docx_file(
        full_path: Path,
        relative_path: str,
        page_start: Optional[int],
        page_end: Optional[int],
        describe: Optional[str],
    ) -> str:
        """Read a Word document with text and visual content."""
        try:
            from docx import Document

            doc = Document(full_path)

            # Extract all text first
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    ]
                    if row_text:
                        text_parts.append(" | ".join(row_text))

            text_content = "\n\n".join(text_parts)

            # Try to get page count and visual content
            try:
                from src.services.document_renderer import get_document_renderer

                renderer = get_document_renderer()
                total_pages = renderer.get_page_count(full_path)

                start = page_start or 1
                end = min(page_end or total_pages, total_pages)

                result_parts = [
                    f"[Pages {start}-{end} of {total_pages}]",
                    "",
                    text_content,
                ]

                # Add visual content for requested pages
                visual_parts = []
                for page_num in range(start, end + 1):
                    visual_content = _get_visual_content(full_path, page_num, describe)
                    if visual_content:
                        visual_parts.append(visual_content)

                if visual_parts:
                    result_parts.append("\n" + "\n".join(visual_parts))

                return "\n".join(result_parts)

            except WorkspaceUnavailableError:
                raise
            except Exception as e:
                # If visual rendering fails, just return text
                logger.debug(f"Could not add visual content for DOCX: {e}")
                return f"[Document: {full_path.name}]\n\n{text_content}"

        except ImportError:
            return "Error: python-docx not installed. Install with: pip install python-docx"
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"DOCX read error for {relative_path}: {e}")
            return f"Error reading Word document: {str(e)}"

    @tool
    def read_file(
        path: str,
        offset: Optional[int] = None,
        limit: Optional[int] = None,
        page_start: Optional[int] = None,
        page_end: Optional[int] = None,
        describe: Optional[str] = None,
    ) -> str:
        """Read content from a file in the workspace.

        For text files, supports line-based access:
        - read_file("doc.txt") - reads first 2000 lines
        - read_file("doc.txt", offset=500, limit=100) - lines 500-599

        For documents (PDF, PPTX, DOCX), supports page-based access:
        - read_file("doc.pdf") - reads first pages within word limit
        - read_file("doc.pdf", page_start=5, page_end=10) - specific pages
        - Visual content (charts, diagrams) is automatically included

        For image files (PNG, JPG, etc.):
        - Returns image data or AI-generated description

        For audio files (MP3, WAV, M4A, OGG, FLAC, etc.):
        - Returns line-numbered text transcription via Whisper
        - Supports offset/limit paging like text files
        - Large files are automatically chunked for transcription

        For archives (ZIP):
        - Returns an entry listing (names, sizes, count) instead of contents

        For any other binary file:
        - Returns a `[binary file: name, N bytes]` message rather than a
          decode error

        Args:
            path: Relative path to the file (e.g., "plan.md")
            offset: For text files: starting line number (1-indexed, default: 1)
            limit: For text files: number of lines to read (default/max: 2000)
            page_start: For documents: first page/slide to read (1-indexed)
            page_end: For documents: last page/slide to read
            describe: Optional query for visual analysis (e.g., "What values are in this chart?")

        Returns:
            File content with line numbers, or error message.
            For documents: includes text + visual content descriptions.
            For images: includes image data or description.
            For audio: includes text transcription of spoken content.
            For archives: an entry listing. For other binaries: a short
            descriptive message.
        """
        try:
            cache_guard_msg = _cloud_cache_guard_for_path(context, path)
            if cache_guard_msg:
                return cache_guard_msg

            # Check file exists
            if not workspace.exists(path):
                return _file_not_found(path)

            full_path = workspace.get_path(path)
            if full_path.is_dir():
                return f"Error: '{path}' is a directory, not a file. Use list_files to see its contents."

            # Handle image files
            # Use local_copy() to ensure the file is on the local
            # filesystem — required for remote workspace backends
            # where get_path() returns a remote-only path.
            if _is_image_file(full_path):
                with workspace.local_copy(path) as local_path:
                    result = _handle_image_file(
                        local_path, describe, display_name=full_path.name
                    )
                if not result.startswith("Error:"):
                    context.record_file_read(path)
                return result

            # Handle audio files (transcribe via Whisper, line-numbered paging)
            if _is_audio_file(full_path):
                with workspace.local_copy(path) as local_path:
                    result = _handle_audio_file(
                        local_path,
                        offset=offset,
                        limit=limit,
                        display_name=full_path.name,
                    )
                if not result.startswith("Error:"):
                    context.record_file_read(path)
                return result

            # Handle visual documents (PDF, PPTX, DOCX) with page-based reading + visual content
            if _is_visual_document(full_path):
                with workspace.local_copy(path) as local_path:
                    result = _read_visual_document(
                        local_path, path, page_start, page_end, describe
                    )
                if not result.startswith("Error:"):
                    context.record_file_read(path)
                return result

            # Handle archives (zip): an entry listing beats a raw codec error
            # and lets the agent ask for a specific member instead of
            # searching the workspace for an "unzip" capability that isn't
            # coming (knowledge-base/knowledge/issues/session_uploads_never_extract_archives.md).
            if _is_archive_file(full_path):
                # Check BEFORE describing the archive: a zip that the
                # upload seam refused to extract (cap/traversal) still
                # parses fine here, so without this it reads as an
                # ordinary, successfully extracted archive — the entries
                # below are shown but none of them are actually separately
                # readable, and nothing else says so.
                note = _read_zip_extraction_note(workspace, path)
                with workspace.local_copy(path) as local_path:
                    archive_size = local_path.stat().st_size
                    result = _describe_zip_archive(
                        local_path, full_path.name, archive_size
                    )
                if note:
                    result = f"{note}\n\n{result}"
                if not result.startswith("Error:"):
                    context.record_file_read(path)
                return result

            # For non-document files, page parameters are ignored
            if page_start is not None or page_end is not None:
                logger.warning(
                    f"page_start/page_end ignored for non-document file: {path}"
                )

            # Apply line-based reading defaults
            start_line = offset if offset is not None else 1
            line_count = limit if limit is not None else DEFAULT_LINE_LIMIT
            line_count = min(line_count, MAX_LINE_LIMIT)  # Cap at max

            if start_line < 1:
                return "Error: offset must be >= 1 (line numbers are 1-indexed)"

            # Read file content. Detect binary content up front: an
            # undecodable file raises UnicodeDecodeError, a ValueError
            # subclass, which the generic handler below would otherwise
            # report as a bare codec message instead of an honest diagnosis.
            try:
                content = workspace.read_file(path)
            except UnicodeDecodeError:
                size = workspace.get_size(path)
                context.record_file_read(path)
                return f"[binary file: {full_path.name}, {size:,} bytes]"
            lines = content.splitlines()
            total_lines = len(lines)

            # Handle empty files: record as read and return informative message
            if total_lines == 0:
                context.record_file_read(path, content)
                return f"File '{path}' is empty (0 lines)."

            # Validate offset
            if start_line > total_lines:
                return (
                    f"Error: offset ({start_line}) exceeds total lines ({total_lines})"
                )

            # Extract requested range (convert to 0-indexed internally)
            end_line = min(start_line + line_count - 1, total_lines)
            selected_lines = lines[start_line - 1 : end_line]

            # Format with line numbers (cat -n style) and truncate long lines
            output_lines = []
            for i, line in enumerate(selected_lines, start=start_line):
                if len(line) > MAX_LINE_LENGTH:
                    line = line[:MAX_LINE_LENGTH] + "..."
                output_lines.append(f"{i:6}\t{line}")

            result = "\n".join(output_lines)

            # Word count cap: dense files (e.g., JSONL) can exceed context budgets
            # even within the 2000-line limit. Truncate to max_read_words if exceeded.
            word_count = len(result.split())
            if word_count > max_read_words:
                words = result.split()
                truncated_text = " ".join(words[:max_read_words])
                # Trim to last complete line to avoid mid-line truncation
                last_newline = truncated_text.rfind("\n")
                if last_newline > 0:
                    truncated_text = truncated_text[:last_newline]
                # Count how many lines we actually kept
                kept_lines = truncated_text.count("\n") + 1
                actual_end = start_line + kept_lines - 1
                result = truncated_text
                result += (
                    f"\n\n[TRUNCATED at word limit ({max_read_words:,} words). "
                    f"Showing lines {start_line}-{actual_end} of {total_lines}. "
                    f"Use offset={actual_end + 1} to continue.]"
                )
            elif end_line < total_lines:
                # Add continuation hint if there are more lines (and no word truncation)
                result += f"\n\n[Lines {start_line}-{end_line} of {total_lines}. "
                result += f"Use offset={end_line + 1} to continue.]"

            # Record successful read for read-before-write tracking
            context.record_file_read(path, content)
            return result

        except FileNotFoundError:
            return _file_not_found(path)
        except ValueError as e:
            return f"Error: {str(e)}"
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"read_file error for {path}: {e}")
            return f"Error reading file: {str(e)}"

    @tool
    def write_file(path: str, content: str) -> str:
        """Write content to a file in the workspace.

        Creates parent directories automatically if they don't exist.
        Overwrites the file if it already exists.

        IMPORTANT: If the file already exists, you must read_file() first to
        understand its current contents before overwriting. This prevents
        accidental data loss from blind overwrites.

        Use this to:
        - Create new files (plan.md, research.md)
        - Save research notes (document_analysis.md)
        - Write intermediate results (candidates/candidates.md)
        - Store processed data (chunks/chunk_001.md)

        Paths are resolved against the **workspace root**, NOT against the
        working directory of your shell. `write_file("output/x.md")` writes the
        same file whether or not `shell_execute` has `cd`'d somewhere — so if
        your shell is in a subdirectory, `cat output/x.md` there will NOT find
        it. The confirmation returns the absolute path; compare it against the
        `CWD:` line in shell results when the two seem to disagree.

        Args:
            path: Path for the file, relative to the workspace root
                (e.g., "research.md")
            content: Content to write

        Returns:
            Confirmation message with the resolved absolute path and size
        """
        # Block binary file extensions — write_file is text-only
        BINARY_EXTENSIONS = {
            ".zip",
            ".tar",
            ".gz",
            ".bz2",
            ".xz",
            ".7z",
            ".rar",
            ".exe",
            ".dll",
            ".so",
            ".dylib",
            ".png",
            ".jpg",
            ".jpeg",
            ".gif",
            ".webp",
            ".bmp",
            ".tiff",
            ".pdf",
            ".docx",
            ".xlsx",
            ".pptx",
            ".db",
            ".sqlite",
            ".pickle",
            ".pkl",
        }
        from pathlib import Path as _P

        if _P(path).suffix.lower() in BINARY_EXTENSIONS:
            return (
                f"Error: Cannot write binary file '{path}'. "
                f"write_file only supports text content. "
                f"Use git tools to commit and push your results for delivery."
            )

        cache_guard_msg = _cloud_cache_guard_for_path(context, path)
        if cache_guard_msg:
            return cache_guard_msg

        upperdir_guard_msg = _cloud_upperdir_guard_for_path(context, path)
        if upperdir_guard_msg:
            return upperdir_guard_msg

        # Enforce word limit
        max_write_words = context.get_config("max_write_words", 10_000)
        word_count = len(content.split())
        if word_count > max_write_words:
            return (
                f"Error: Content too large ({word_count:,} words, limit is {max_write_words:,}). "
                f"Split into multiple write_file calls — for example, write each section or chapter "
                f"separately using different file paths, then combine if needed."
            )

        try:
            # Enforce read-before-write for existing non-empty files. A
            # versioned read must still match the current complete text; this
            # catches Canvas/user edits even when a best-effort invalidation was
            # missed while the agent was detached.
            existed = workspace.exists(path)
            had_recent_read = context.was_recently_read(path)
            if existed:
                existing = workspace.read_file(path)
                if existing.strip() and not context.recent_read_matches(path, existing):
                    context.invalidate_recent_read(path)
                    from src.services.guardrails import format_nudge

                    model = (
                        context._llm_config.model
                        if context._llm_config is not None
                        else None
                    )
                    return format_nudge(
                        "read_file_required_error",
                        model=model,
                        file_path=path,
                        tool_name="write_file",
                    )

            # Snapshot for undo before writing
            if context._snapshot_callback:
                context._snapshot_callback(path)

            workspace.write_file(path, content)
            if existed and had_recent_read:
                context.record_file_read(path, content)

            # Absolute, not the caller's own string: `output/x.md` here and
            # `output/x.md` in a shell that has cd'd elsewhere are different
            # files, and only this line says which one you got.
            return f"Written: {_absolute(workspace, path)} ({len(content)} chars)"

        except ValueError as e:
            return f"Error: {str(e)}"
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"write_file error for {path}: {e}")
            return f"Error writing file: {str(e)}"

    @tool
    def edit_file(
        path: str,
        old_string: str = "",
        new_string: str = "",
        position: Optional[Literal["start", "end"]] = None,
    ) -> str:
        """Edit a file by replacing text or inserting at start/end.

        IMPORTANT: You must read_file() before editing. This ensures you
        understand the file's current contents before modifying it.

        **Modes:**

        1. **Replace mode** (default): Set `old_string` and `new_string` to find
           and replace text. The `old_string` must appear exactly once.

        2. **Append mode**: Set `position="end"` to add `new_string` at the end
           of the file. The `old_string` parameter is ignored.

        3. **Prepend mode**: Set `position="start"` to add `new_string` at the
           beginning of the file. The `old_string` parameter is ignored.

        Args:
            path: Relative path to the file (e.g., "plan.md")
            old_string: Text to find and replace (required for replace mode)
            new_string: Replacement text or content to insert
            position: Insert position - "start", "end", or None for replace mode

        Returns:
            Confirmation message or error with guidance
        """
        try:
            cache_guard_msg = _cloud_cache_guard_for_path(context, path)
            if cache_guard_msg:
                return cache_guard_msg

            upperdir_guard_msg = _cloud_upperdir_guard_for_path(context, path)
            if upperdir_guard_msg:
                return upperdir_guard_msg

            if not workspace.exists(path):
                return _file_not_found(path)

            full_path = workspace.get_path(path)
            if full_path.is_dir():
                return f"Error: '{path}' is a directory, not a file."

            # Enforce read-before-write discipline (skip for empty files), and
            # reject a stale versioned read after an out-of-band edit.
            content = workspace.read_file(path)
            had_recent_read = context.was_recently_read(path)
            if content.strip() and not context.recent_read_matches(path, content):
                context.invalidate_recent_read(path)
                from src.services.guardrails import format_nudge

                model = (
                    context._llm_config.model
                    if context._llm_config is not None
                    else None
                )
                return format_nudge(
                    "read_file_required_error",
                    model=model,
                    file_path=path,
                    tool_name="edit_file",
                )

            # Validate position parameter
            if position is not None and position not in ("start", "end"):
                return (
                    f"Error: Invalid position '{position}'. "
                    f"Use 'start' to prepend, 'end' to append, or omit for replace mode."
                )

            # Snapshot for undo before editing
            if context._snapshot_callback:
                context._snapshot_callback(path)

            # Position-based insert modes
            if position == "end":
                new_content = content + new_string
                workspace.write_file(path, new_content)
                if had_recent_read:
                    context.record_file_read(path, new_content)
                return f"Appended to: {path}"

            if position == "start":
                new_content = new_string + content
                workspace.write_file(path, new_content)
                if had_recent_read:
                    context.record_file_read(path, new_content)
                return f"Prepended to: {path}"

            # Replace mode (default) - requires old_string
            if not old_string:
                return (
                    "Error: old_string is required for replace mode. "
                    "To append, use position='end'. To prepend, use position='start'."
                )

            count = content.count(old_string)

            if count == 0:
                # Show a short snippet of the file to help the caller orient
                preview = content[:200].replace("\n", "\\n")
                return (
                    f"Error: old_string not found in {path}. "
                    f"Make sure the string matches exactly (including whitespace and newlines). "
                    f"File starts with: {preview!r}"
                )

            if count > 1:
                return (
                    f"Error: old_string appears {count} times in {path}. "
                    f"Include more surrounding context to make the match unique."
                )

            new_content = content.replace(old_string, new_string, 1)
            workspace.write_file(path, new_content)
            if had_recent_read:
                context.record_file_read(path, new_content)

            return f"Edited: {path}"

        except ValueError as e:
            return f"Error: {str(e)}"
        except WorkspaceUnavailableError:
            raise
        except Exception as e:
            logger.error(f"edit_file error for {path}: {e}")
            return f"Error editing file: {str(e)}"

    return [read_file, write_file, edit_file]
